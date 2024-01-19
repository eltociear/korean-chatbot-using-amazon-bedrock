import json
import boto3
import os
import traceback
import PyPDF2
import time
import docx

from io import BytesIO
from urllib import parse
from botocore.config import Config
from urllib.parse import unquote_plus
from langchain_community.embeddings import BedrockEmbeddings
from langchain.vectorstores.opensearch_vector_search import OpenSearchVectorSearch
from langchain.docstore.document import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from opensearchpy import OpenSearch
from pptx import Presentation

s3 = boto3.client('s3')
s3_bucket = os.environ.get('s3_bucket') # bucket name
s3_prefix = os.environ.get('s3_prefix')
meta_prefix = "metadata/"
kendra_region = os.environ.get('kendra_region', 'us-west-2')

opensearch_account = os.environ.get('opensearch_account')
opensearch_passwd = os.environ.get('opensearch_passwd')
opensearch_url = os.environ.get('opensearch_url')
kendraIndex = os.environ.get('kendraIndex')
sqsUrl = os.environ.get('sqsUrl')
sqs = boto3.client('sqs')

roleArn = os.environ.get('roleArn') 
path = os.environ.get('path')
max_object_size = int(os.environ.get('max_object_size'))

capabilities = json.loads(os.environ.get('capabilities'))
print('capabilities: ', capabilities)

os_client = OpenSearch(
    hosts = [{
        'host': opensearch_url.replace("https://", ""), 
        'port': 443
    }],
    http_compress = True,
    http_auth=(opensearch_account, opensearch_passwd),
    use_ssl = True,
    verify_certs = True,
    ssl_assert_hostname = False,
    ssl_show_warn = False,
)

def delete_index_if_exist(index_name):    
    if os_client.indices.exists(index_name):
        print('delete opensearch document index: ', index_name)
        response = os_client.indices.delete(
            index=index_name
        )
        print('removed index: ', response)    
    else:
        print('no index: ', index_name)

# Kendra
kendra_client = boto3.client(
    service_name='kendra', 
    region_name=kendra_region,
    config = Config(
        retries=dict(
            max_attempts=10
        )
    )
)

# embedding for RAG
bedrock_region = "us-west-2"
    
boto3_bedrock = boto3.client(
    service_name='bedrock-runtime',
    region_name=bedrock_region,
    config=Config(
        retries = {
            'max_attempts': 30
        }            
    )
)

bedrock_embeddings = BedrockEmbeddings(
    client=boto3_bedrock,
    region_name = bedrock_region,
    model_id = 'amazon.titan-embed-text-v1' 
)   

def store_document_for_opensearch(bedrock_embeddings, docs, documentId):
    index_name = "rag-index-"+documentId
    print('index_name: ', index_name)
    
    if len(index_name)>=255:
        index_name = index_name[1:255]
        print('index_name: ', index_name)
    
    delete_index_if_exist(index_name)

    try:
        vectorstore = OpenSearchVectorSearch(
            index_name=index_name,  
            is_aoss = False,
            #engine="faiss",  # default: nmslib
            embedding_function = bedrock_embeddings,
            opensearch_url = opensearch_url,
            http_auth=(opensearch_account, opensearch_passwd),
        )
        response = vectorstore.add_documents(docs, bulk_size = 2000)
        print('response of adding documents: ', response)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                
        #raise Exception ("Not able to request to LLM")

    print('uploaded into opensearch')
    
def store_document_for_opensearch_with_nori(bedrock_embeddings, docs, documentId):
    index_name = "rag-index-"+documentId
    print('index_name: ', index_name)
    
    if len(index_name)>=255:
        index_name = index_name[1:255]
        print('index_name: ', index_name)
    
    delete_index_if_exist(index_name)
    
    index_body = {
        'settings': {
            'analysis': {
                'analyzer': {
                    'my_analyzer': {
                        'char_filter': ['html_strip'], 
                        'tokenizer': 'nori',
                        'filter': ['nori_number','lowercase','trim','my_nori_part_of_speech'],
                        'type': 'custom'
                    }
                },
                'tokenizer': {
                    'nori': {
                        'decompound_mode': 'mixed',
                        'discard_punctuation': 'true',
                        'type': 'nori_tokenizer'}
                },
                "filter": {
                    "my_nori_part_of_speech": {
                        "type": "nori_part_of_speech",
                        "stoptags": [
                                "E", "IC", "J", "MAG", "MAJ",
                                "MM", "SP", "SSC", "SSO", "SC",
                                "SE", "XPN", "XSA", "XSN", "XSV",
                                "UNA", "NA", "VSV"
                        ]
                    }
                }
            },
            'index': {
                'knn': True,
                'knn.space_type': 'cosinesimil'  # Example space type
            }
        },
        'mappings': {
            'properties': {
                'metadata': {
                    'properties': {
                        'source' : {'type': 'keyword'},                    
                        'last_updated': {'type': 'date'},
                        'project': {'type': 'keyword'},
                        'seq_num': {'type': 'long'},
                        'title': {'type': 'text'},  # For full-text search
                        'url': {'type': 'text'},  # For full-text search
                    }
                },            
                'text': {
                    'analyzer': 'my_analyzer',
                    'search_analyzer': 'my_analyzer',
                    'type': 'text'
                },
                'vector_field': {
                    'type': 'knn_vector',
                    'dimension': 1536  # Replace with your vector dimension
                }
            }
        }
    }
    
    try: # create index
        response = os_client.indices.create(
            index_name,
            body=index_body
        )
        print('\nCreating index:')
        print(response)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                
        #raise Exception ("Not able to create the index")

    try: # 
        vectorstore = OpenSearchVectorSearch(
            index_name=index_name,  
            is_aoss = False,
            #engine="faiss",  # default: nmslib
            embedding_function = bedrock_embeddings,
            opensearch_url = opensearch_url,
            http_auth=(opensearch_account, opensearch_passwd),
        )
        response = vectorstore.add_documents(docs, bulk_size = 2000)
        print('response of adding documents: ', response)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                
        #raise Exception ("Not able to request to LLM")

    print('uploaded into opensearch')    
 
# store document into Kendra
def store_document_for_kendra(path, key, documentId):
    print('store document into kendra')
    encoded_name = parse.quote(key)
    source_uri = path+encoded_name    
    #print('source_uri: ', source_uri)
    ext = (key[key.rfind('.')+1:len(key)]).upper()
    print('ext: ', ext)

    # PLAIN_TEXT, XSLT, MS_WORD, RTF, CSV, JSON, HTML, PDF, PPT, MD, XML, MS_EXCEL
    if(ext == 'PPTX'):
        file_type = 'PPT'
    elif(ext == 'TXT'):
        file_type = 'PLAIN_TEXT'         
    elif(ext == 'XLS' or ext == 'XLSX'):
        file_type = 'MS_EXCEL'      
    elif(ext == 'DOC' or ext == 'DOCX'):
        file_type = 'MS_WORD'
    else:
        file_type = ext

    kendra_client = boto3.client(
        service_name='kendra', 
        region_name=kendra_region,
        config = Config(
            retries=dict(
                max_attempts=10
            )
        )
    )

    documents = [
        {
            "Id": documentId,
            "Title": key,
            "S3Path": {
                "Bucket": s3_bucket,
                "Key": key
            },
            "Attributes": [
                {
                    "Key": '_source_uri',
                    'Value': {
                        'StringValue': source_uri
                    }
                },
                {
                    "Key": '_language_code',
                    'Value': {
                        'StringValue': "ko"
                    }
                },
            ],
            "ContentType": file_type
        }
    ]
    print('document info: ', documents)

    try:
        result = kendra_client.batch_put_document(
            IndexId = kendraIndex,
            RoleArn = roleArn,
            Documents = documents       
        )
        print('batch_put_document(kendra): ', result)
        print('uploaded into kendra')
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
        # raise Exception ("Not able to put a document in Kendra")

def create_metadata(bucket, key, meta_prefix, s3_prefix, uri, category, documentId):    
    title = key
    timestamp = int(time.time())

    metadata = {
        "Attributes": {
            "_category": category,
            "_source_uri": uri,
            "_version": str(timestamp),
            "_language_code": "ko"
        },
        "Title": title,
        "DocumentId": documentId,        
    }
    print('metadata: ', metadata)
    
    objectName = (key[key.find(s3_prefix)+len(s3_prefix)+1:len(key)]).upper()
    print('objectName: ', objectName)

    client = boto3.client('s3')
    try: 
        client.put_object(
            Body=json.dumps(metadata), 
            Bucket=bucket, 
            Key=meta_prefix+objectName+'.metadata.json' 
        )
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
        raise Exception ("Not able to create meta file")

# load documents from s3 for pdf and txt
def load_document(file_type, key):
    s3r = boto3.resource("s3")
    doc = s3r.Object(s3_bucket, key)
    
    contents = ""
    if file_type == 'pdf':
        Byte_contents = doc.get()['Body'].read()
        
        try: 
            reader = PyPDF2.PdfReader(BytesIO(Byte_contents))
            
            texts = []
            for page in reader.pages:
                texts.append(page.extract_text())
            contents = '\n'.join(texts)
        except Exception:
                err_msg = traceback.format_exc()
                print('err_msg: ', err_msg)
                # raise Exception ("Not able to load the pdf file")
                
    elif file_type == 'pptx':
        Byte_contents = doc.get()['Body'].read()
            
        try:
            prs = Presentation(BytesIO(Byte_contents))

            texts = []
            for i, slide in enumerate(prs.slides):
                text = ""
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = text + shape.text
                texts.append(text)
            contents = '\n'.join(texts)
        except Exception:
                err_msg = traceback.format_exc()
                print('err_msg: ', err_msg)
                # raise Exception ("Not able to load texts from preseation file")
        
    elif file_type == 'txt':       
        try:  
            contents = doc.get()['Body'].read().decode('utf-8')
        except Exception:
            err_msg = traceback.format_exc()
            print('error message: ', err_msg)        
            # raise Exception ("Not able to load the file")

    elif file_type == 'docx':
        try:
            Byte_contents = doc.get()['Body'].read()                    
            doc_contents =docx.Document(BytesIO(Byte_contents))

            texts = []
            for i, para in enumerate(doc_contents.paragraphs):
                if(para.text):
                    texts.append(para.text)
                    # print(f"{i}: {para.text}")        
            contents = '\n'.join(texts)            
            # print('contents: ', contents)
        except Exception:
                err_msg = traceback.format_exc()
                print('err_msg: ', err_msg)
                # raise Exception ("Not able to load docx")   
    
    texts = ""
    if len(contents)>0:
        new_contents = str(contents).replace("\n"," ") 
        print('length: ', len(new_contents))
        
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=100,
            separators=["\n\n", "\n", ".", " ", ""],
            length_function = len,
        ) 

        texts = text_splitter.split_text(new_contents) 
                        
    return texts

def check_supported_type(file_type, size):
    if size > 5000 and size<max_object_size and (file_type == 'pdf' or file_type == 'txt' or file_type == 'csv' or file_type == 'pptx' or file_type == 'ppt' or file_type == 'docx' or file_type == 'doc' or file_type == 'xlsx'):
        return True
    if size > 0 and file_type == 'txt':
        return True
    else:
        return False
    
# load csv documents from s3
def lambda_handler(event, context):
    print('event: ', event)    
    
    documentIds = []
    for record in event['Records']:
        receiptHandle = record['receiptHandle']
        print("receiptHandle: ", receiptHandle)
        
        body = record['body']
        print("body: ", body)
        
        jsonbody = json.loads(body)        
        bucket = jsonbody['bucket']        
        # translate utf8
        #key = unquote_plus(jsonbody['key']) # url decoding
        key = jsonbody['key'] # url decoding
        print('bucket: ', bucket)
        print('key: ', key)        
        eventName = jsonbody['type']
        
        start_time = time.time()      
        
        file_type = key[key.rfind('.')+1:len(key)].lower()
        print('file_type: ', file_type)
            
        size = 0
        try:
            result = s3.get_object_attributes(Bucket=bucket, Key=key, ObjectAttributes=['ObjectSize'])  
            print('result: ', result)
            
            size = int(result['ObjectSize'])
            print('object size: ', size)
        except Exception:
            err_msg = traceback.format_exc()
            print('err_msg: ', err_msg)
            # raise Exception ("Not able to get object info") 
        
        if eventName == 'ObjectRemoved:Delete':
            if check_supported_type(file_type, size):
                objectName = (key[key.find(s3_prefix)+len(s3_prefix)+1:len(key)]).upper()
                print('objectName: ', objectName)
                
                # get metadata from s3
                metadata_key = meta_prefix+objectName+'.metadata.json'
                print('metadata_key: ', metadata_key)

                documentId = ""
                try: 
                    metadata_obj = s3.get_object(Bucket=bucket, Key=metadata_key)
                    metadata_body = metadata_obj['Body'].read().decode('utf-8')
                    metadata = json.loads(metadata_body)
                    print('metadata: ', metadata)
                    documentId = metadata['DocumentId']
                    print('documentId: ', documentId)
                    documentIds.append(documentId)
                except Exception:
                    err_msg = traceback.format_exc()
                    print('err_msg: ', err_msg)
                    # raise Exception ("Not able to get the object")
                    
                if documentId:
                    try: # delete metadata                        
                        print('delete metadata: ', metadata_key)                        
                        result = s3.delete_object(Bucket=bucket, Key=metadata_key)
                        # print('result of metadata deletion: ', result)
                        
                        # delete document index of opensearch
                        index_name = "rag-index-"+documentId
                        # print('index_name: ', index_name)
                        delete_index_if_exist(index_name)                    
                    except Exception:
                        err_msg = traceback.format_exc()
                        print('err_msg: ', err_msg)
                        # raise Exception ("Not able to delete documents in Kendra")
                    
                    # delete kendra documents
                    print('delete kendra documents: ', documentIds)            
                    try: 
                        result = kendra_client.batch_delete_document(
                            IndexId = kendraIndex,
                            DocumentIdList=[
                                documentId,
                            ]
                        )
                        print('result: ', result)
                    except Exception:
                        err_msg = traceback.format_exc()
                        print('err_msg: ', err_msg)
                        # raise Exception ("Not able to delete documents in Kendra")
                    
        elif eventName == "ObjectCreated:Put":
            category = "upload"
            documentId = category + "-" + key
            documentId = documentId.replace(' ', '_') # remove spaces
            documentId = documentId.replace(',', '_') # remove commas
            documentId = documentId.replace('/', '_') # remove slash
            documentId = documentId.lower() # change to lowercase
            print('documentId: ', documentId)
            
            if check_supported_type(file_type, size): 
                for type in capabilities:                
                    if type=='kendra':         
                        print('upload to kendra: ', key)                                                
                        # PLAIN_TEXT, XSLT, MS_WORD, RTF, CSV, JSON, HTML, PDF, PPT, MD, XML, MS_EXCEL                    
                        store_document_for_kendra(path, key, documentId)  # store the object into kendra
                        
                    elif type=='opensearch':
                        if file_type == 'pdf' or file_type == 'txt' or file_type == 'csv' or file_type == 'pptx' or file_type == 'docx':
                            print('upload to opensearch: ', key) 
                            texts = load_document(file_type, key)
                            
                            docs = []
                            for i in range(len(texts)):
                                if texts[i]:
                                    docs.append(
                                        Document(
                                            page_content=texts[i],
                                            metadata={
                                                'name': key,
                                                # 'page':i+1,
                                                'uri': path+parse.quote(key)
                                            }
                                        )
                                    )
                            print('docs size: ', len(docs))
                            if len(docs)>0:
                                print('docs[0]: ', docs[0])                            
                                # store_document_for_opensearch(bedrock_embeddings, docs, documentId)
                                store_document_for_opensearch_with_nori(bedrock_embeddings, docs, documentId)
                    
                create_metadata(bucket=s3_bucket, key=key, meta_prefix=meta_prefix, s3_prefix=s3_prefix, uri=path+parse.quote(key), category=category, documentId=documentId)
            else: # delete if the object is unsupported one for format or size
                try:
                    print('delete unsupported file: ', key)                                
                    result = s3.delete_object(Bucket=bucket, Key=key)
                    print('result of deletion of the unsupported file: ', result)
                            
                except Exception:
                    err_msg = traceback.format_exc()
                    print('err_msg: ', err_msg)
                    # raise Exception ("Not able to delete unsupported file")

        print('processing time: ', str(time.time() - start_time))
        
        # delete queue
        try:
            sqs.delete_message(QueueUrl=sqsUrl, ReceiptHandle=receiptHandle)
        except Exception as e:        
            print('Fail to delete the queue message: ', e)
    return {
        'statusCode': 200
    }
